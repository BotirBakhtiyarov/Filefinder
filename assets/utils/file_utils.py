import os
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
import PyPDF2

def scan_files(directory, extensions):
    return [
        os.path.join(root, f)
        for root, _, files in os.walk(directory)
        for f in files
        if f.split(".")[-1].lower() in extensions
    ]

def extract_text_content(file_path):
    ext = file_path.split(".")[-1].lower()
    try:
        if ext == "docx":
            return "\n".join(p.text for p in Document(file_path).paragraphs)
        elif ext == "xlsx":
            return " ".join(str(cell) for sheet in load_workbook(file_path) for row in sheet.values for cell in row)
        elif ext == "pptx":
            return "\n".join(shape.text for slide in Presentation(file_path).slides for shape in slide.shapes if shape.has_text_frame)
        elif ext == "pdf":
            pdf_reader = PyPDF2.PdfReader(file_path)
            return "\n".join(page.extract_text() or "" for page in pdf_reader.pages)
        elif ext == "txt":
            with open(file_path, "r", encoding="utf-8") as f:
                return f.read()
    except Exception as e:
        print(f"Error reading {file_path}: {str(e)}")
    return ""